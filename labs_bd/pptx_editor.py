# Import the Presentation class from the pptx module.
from tkinter import Image
from pptx import Presentation

print('Напишите имя пользователя')
username = input()
prof = []
points = []
str_prof = ""
while True:
    print('Напишите название профессии и процент набранных баллов за неё, если закончили, наберите "exit"')
    print('Название профессии')
    prof_add = input()
    prof.append(prof_add)
    if 'exit' in prof:
        break
    print('Кол-во процентов')
    points_add = input()
    points.append(points_add)
    if 'exit' in points:
        break

# Create an instance of the Presentation class.
prs = Presentation()
# Create the title slide.
title_slide_layout = prs.slide_layouts[0]
# Add the title slide to the PPT slides array.
slide = prs.slides.add_slide(title_slide_layout)

slide.background.fill.background()
slide.background.fill.background_image = "labs_bd/bg.jpg"
# Get the PPT title object.
title = slide.shapes.title
# Get the PPT subtitle object.
subtitle = slide.placeholders[1]

# Set the PPT title.
title.text = f"Сертификат выдан пользователю {username}"
# Set the PPT title slide sub title.
for i in range(len(prof)):
    str_prof += "".join(prof[i]) + ": ".join(points[i] + "\n")

subtitle.text = str_prof
# for prof, points in score:
#     subtitle.text = f"Профессия: {prof}, набранный процент: {points}"

# Save the PPT to a .pptx file.
prs.save("labs_bd/test.pptx")